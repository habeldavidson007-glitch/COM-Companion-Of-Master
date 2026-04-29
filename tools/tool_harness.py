"""
COM Companion Tool Harness
==========================
Bridges LLM input to Python tools with 6 key optimizations:
1. Tool Pre-validation & Health Check
2. Payload Validation Before Execution
3. Tool Execution Caching
4. Parallel Signal Execution
5. Tool-Specific Error Recovery
6. Output File Path Management

This module is isolated from the core system and focuses solely on tool execution.
"""

import os
import re
import time
import hashlib
import logging
from datetime import datetime
from typing import Dict, List, Optional, Any, Tuple
from concurrent.futures import ThreadPoolExecutor, as_completed
from collections import OrderedDict
import threading

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# =============================================================================
# CONFIGURATION
# =============================================================================

# Output directory for all generated files
OUTPUT_DIR = os.environ.get("COM_OUTPUT_DIR", "./com_output")

# Cache settings
CACHE_MAX_SIZE = 100
CACHE_ENABLED = True

# Retry settings
MAX_RETRIES = 3
RETRY_DELAY_BASE = 1.0  # seconds

# =============================================================================
# 6. OUTPUT FILE PATH MANAGEMENT
# =============================================================================

class FilePathManager:
    """Manages output file paths with sanitization, collision avoidance, and directory creation."""
    
    def __init__(self, base_dir: str = OUTPUT_DIR):
        self.base_dir = os.path.abspath(base_dir)
        self._lock = threading.Lock()
        self._ensure_directory_exists()
    
    def _ensure_directory_exists(self):
        """Create output directory if it doesn't exist."""
        if not os.path.exists(self.base_dir):
            os.makedirs(self.base_dir, exist_ok=True)
            logger.info(f"Created output directory: {self.base_dir}")
    
    def sanitize_filename(self, filename: str) -> str:
        """Remove invalid characters from filename."""
        # Remove or replace invalid characters
        invalid_chars = r'<>:"/\\|?*'
        sanitized = re.sub(f'[{re.escape(invalid_chars)}]', '_', filename)
        # Remove leading/trailing spaces and dots
        sanitized = sanitized.strip(' .')
        # Limit length
        if len(sanitized) > 200:
            name, ext = os.path.splitext(sanitized)
            sanitized = name[:200-len(ext)] + ext
        return sanitized or "unnamed_file"
    
    def get_unique_path(self, filename: str, extension: str) -> str:
        """Generate a unique file path, avoiding collisions with timestamps.
        
        Uses atomic file creation (O_CREAT | O_EXCL) to prevent TOCTOU race
        conditions when multiple threads request the same path concurrently.
        """
        sanitized_name = self.sanitize_filename(filename)
        if not sanitized_name.endswith(extension):
            base_name = sanitized_name
        else:
            base_name = sanitized_name[:-len(extension)]
        
        # First attempt: original name - atomically claim it
        candidate_path = os.path.join(self.base_dir, f"{base_name}{extension}")
        
        with self._lock:
            try:
                # Exclusive create — atomically claims the path
                fd = os.open(candidate_path, os.O_CREAT | os.O_EXCL | os.O_WRONLY)
                os.close(fd)
                return candidate_path
            except FileExistsError:
                pass
            
            # Collision detected: add timestamp
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S_%f")
            candidate_path = os.path.join(self.base_dir, f"{base_name}_{timestamp}{extension}")
            
            # Atomically create the timestamped path
            fd = os.open(candidate_path, os.O_CREAT | os.O_EXCL | os.O_WRONLY)
            os.close(fd)
            
            return candidate_path
    
    def resolve_path(self, relative_path: str) -> str:
        """Convert relative path to absolute path."""
        if os.path.isabs(relative_path):
            return relative_path
        return os.path.abspath(os.path.join(self.base_dir, relative_path))

# Global file path manager
file_manager = FilePathManager(OUTPUT_DIR)

# =============================================================================
# 3. TOOL EXECUTION CACHING
# =============================================================================

class LRUCache:
    """Thread-safe LRU cache for tool execution results."""
    
    def __init__(self, max_size: int = CACHE_MAX_SIZE):
        self.max_size = max_size
        self.cache: OrderedDict[str, Any] = OrderedDict()
        self.lock = threading.Lock()
        self.hits = 0
        self.misses = 0
    
    def _generate_key(self, tool_type: str, payload: str) -> str:
        """Generate a unique cache key using SHA256."""
        content = f"{tool_type}:{payload}"
        return hashlib.sha256(content.encode()).hexdigest()
    
    def get(self, tool_type: str, payload: str) -> Optional[Any]:
        """Get cached result if exists and valid."""
        if not CACHE_ENABLED:
            return None
        
        key = self._generate_key(tool_type, payload)
        
        with self.lock:
            if key in self.cache:
                result = self.cache[key]
                # Validate file exists for file-based results
                if isinstance(result, dict) and 'file_path' in result:
                    if not os.path.exists(result['file_path']):
                        # File was deleted, invalidate cache
                        del self.cache[key]
                        self.misses += 1
                        return None
                
                # Move to end (most recently used)
                self.cache.move_to_end(key)
                self.hits += 1
                return result
            
            self.misses += 1
            return None
    
    def set(self, tool_type: str, payload: str, result: Any):
        """Store result in cache."""
        if not CACHE_ENABLED:
            return
        
        key = self._generate_key(tool_type, payload)
        
        with self.lock:
            if key in self.cache:
                self.cache.move_to_end(key)
            self.cache[key] = result
            
            # Evict oldest if over capacity
            while len(self.cache) > self.max_size:
                self.cache.popitem(last=False)
    
    def clear(self):
        """Clear all cached results."""
        with self.lock:
            self.cache.clear()
            self.hits = 0
            self.misses = 0
    
    def get_stats(self) -> Dict[str, int]:
        """Return cache statistics."""
        with self.lock:
            total = self.hits + self.misses
            hit_rate = (self.hits / total * 100) if total > 0 else 0
            return {
                'size': len(self.cache),
                'max_size': self.max_size,
                'hits': self.hits,
                'misses': self.misses,
                'hit_rate': round(hit_rate, 2)
            }

# Global cache instance
tool_cache = LRUCache(CACHE_MAX_SIZE)

# =============================================================================
# 1. TOOL PRE-VALIDATION & HEALTH CHECK
# =============================================================================

class ToolHealthChecker:
    """Validates tool availability and dependencies before execution."""
    
    def __init__(self):
        self.tool_status: Dict[str, Dict[str, Any]] = {}
        self._check_all_tools()
    
    def _check_all_tools(self):
        """Check health of all available tools."""
        self.tool_status = {
            'XLS': self._check_excel_tool(),
            'PPT': self._check_powerpoint_tool(),
            'PDF': self._check_pdf_tool(),
            'GODOT': self._check_godot_tool()
        }
    
    def _check_excel_tool(self) -> Dict[str, Any]:
        """Check if Excel tool dependencies are available."""
        status = {'available': False, 'reason': '', 'dependencies': []}
        
        try:
            import pandas
            status['dependencies'].append('pandas: OK')
        except ImportError:
            status['reason'] = 'pandas not installed'
            return status
        
        try:
            import openpyxl
            status['dependencies'].append('openpyxl: OK')
        except ImportError:
            status['reason'] = 'openpyxl not installed'
            return status
        
        status['available'] = True
        return status
    
    def _check_powerpoint_tool(self) -> Dict[str, Any]:
        """Check if PowerPoint tool dependencies are available."""
        status = {'available': False, 'reason': '', 'dependencies': []}
        
        try:
            from pptx import Presentation
            status['dependencies'].append('python-pptx: OK')
        except ImportError:
            status['reason'] = 'python-pptx not installed'
            return status
        
        status['available'] = True
        return status
    
    def _check_pdf_tool(self) -> Dict[str, Any]:
        """Check if PDF tool dependencies are available."""
        status = {'available': False, 'reason': '', 'dependencies': []}
        
        try:
            from fpdf import FPDF
            status['dependencies'].append('fpdf: OK')
        except ImportError:
            status['reason'] = 'fpdf not installed (pip install fpdf)'
            return status
        
        status['available'] = True
        return status
    
    def _check_godot_tool(self) -> Dict[str, Any]:
        """Check if Godot tool requirements are met."""
        status = {'available': False, 'reason': '', 'dependencies': []}
        
        # Godot tool typically just needs filesystem access
        if os.access(OUTPUT_DIR, os.W_OK):
            status['dependencies'].append('filesystem: OK')
            status['available'] = True
        else:
            status['reason'] = 'No write permission to output directory'
        
        return status
    
    def is_tool_available(self, tool_type: str) -> bool:
        """Check if a specific tool is available."""
        tool_type = tool_type.upper()
        if tool_type not in self.tool_status:
            return False
        return self.tool_status[tool_type].get('available', False)
    
    def get_unavailable_tools(self) -> List[str]:
        """Return list of unavailable tools."""
        return [
            tool for tool, status in self.tool_status.items()
            if not status.get('available', False)
        ]
    
    def get_health_report(self) -> str:
        """Generate a human-readable health report."""
        lines = ["=== Tool Health Report ==="]
        for tool, status in self.tool_status.items():
            avail = "✅" if status['available'] else "❌"
            lines.append(f"{avail} {tool}: {'Available' if status['available'] else status['reason']}")
            if status['dependencies']:
                for dep in status['dependencies']:
                    lines.append(f"   └─ {dep}")
        
        cache_stats = tool_cache.get_stats()
        lines.append(f"\n📊 Cache: {cache_stats['size']}/{cache_stats['max_size']} entries, {cache_stats['hit_rate']}% hit rate")
        
        return "\n".join(lines)

# Global health checker
health_checker = ToolHealthChecker()

def validate_tool_health() -> bool:
    """Validate all tools are healthy before use."""
    unavailable = health_checker.get_unavailable_tools()
    if unavailable:
        logger.warning(f"Unavailable tools: {', '.join(unavailable)}")
        return False
    return True

# =============================================================================
# 2. PAYLOAD VALIDATION BEFORE EXECUTION
# =============================================================================

class PayloadValidator:
    """Validates signal payloads before execution to prevent runtime errors."""
    
    @staticmethod
    def validate_xls_payload(payload: str) -> Tuple[bool, str]:
        """
        Validate XLS payload format: filename:col1,col2,col3
        Returns (is_valid, error_message)
        """
        parts = payload.split(':')
        if len(parts) != 2:
            return False, "XLS payload must be 'filename:col1,col2,col3'"
        
        filename, columns = parts
        if not filename.strip():
            return False, "Filename cannot be empty"
        
        if not columns.strip():
            return False, "Columns cannot be empty"
        
        col_list = [c.strip() for c in columns.split(',') if c.strip()]
        if not col_list:
            return False, "At least one column must be specified"
        
        return True, ""
    
    @staticmethod
    def validate_ppt_payload(payload: str) -> Tuple[bool, str]:
        """
        Validate PPT payload format: filename:slide1|slide2|slide3
        Returns (is_valid, error_message)
        """
        parts = payload.split(':')
        if len(parts) != 2:
            return False, "PPT payload must be 'filename:slide1|slide2|slide3'"
        
        filename, slides = parts
        if not filename.strip():
            return False, "Filename cannot be empty"
        
        if not slides.strip():
            return False, "Slides cannot be empty"
        
        slide_list = [s.strip() for s in slides.split('|') if s.strip()]
        if not slide_list:
            return False, "At least one slide must be specified"
        
        return True, ""
    
    @staticmethod
    def validate_pdf_payload(payload: str) -> Tuple[bool, str]:
        """
        Validate PDF payload format: filename:content
        Returns (is_valid, error_message)
        """
        parts = payload.split(':', 1)  # Split only on first colon
        if len(parts) != 2:
            return False, "PDF payload must be 'filename:content'"
        
        filename, content = parts
        if not filename.strip():
            return False, "Filename cannot be empty"
        
        if not content.strip():
            return False, "Content cannot be empty"
        
        return True, ""
    
    @staticmethod
    def validate_godot_payload(payload: str) -> Tuple[bool, str]:
        """
        Validate GODOT payload format: template_name:config_json_or_params
        Returns (is_valid, error_message)
        """
        parts = payload.split(':')
        if len(parts) < 2:
            return False, "GODOT payload must be 'template_name:params'"
        
        template = parts[0]
        if not template.strip():
            return False, "Template name cannot be empty"
        
        return True, ""
    
    @classmethod
    def validate_payload(cls, tool_type: str, payload: str) -> Tuple[bool, str]:
        """Validate payload based on tool type."""
        tool_type = tool_type.upper()
        
        validators = {
            'XLS': cls.validate_xls_payload,
            'PPT': cls.validate_ppt_payload,
            'PDF': cls.validate_pdf_payload,
            'GODOT': cls.validate_godot_payload
        }
        
        validator = validators.get(tool_type)
        if not validator:
            return False, f"Unknown tool type: {tool_type}"
        
        return validator(payload)

# =============================================================================
# 5. TOOL-SPECIFIC ERROR RECOVERY
# =============================================================================

class ErrorRecoveryHandler:
    """Handles tool-specific error recovery with retry logic."""
    
    @staticmethod
    def should_retry(exception: Exception, attempt: int) -> bool:
        """Determine if an error should be retried."""
        if attempt >= MAX_RETRIES:
            return False
        
        # Don't retry import errors
        if isinstance(exception, ImportError):
            return False
        
        # Retry on transient errors
        retry_exceptions = (PermissionError, OSError, ConnectionRefusedError, TimeoutError)
        return isinstance(exception, retry_exceptions)
    
    @staticmethod
    def get_delay(attempt: int) -> float:
        """Calculate exponential backoff delay."""
        return RETRY_DELAY_BASE * (2 ** attempt)
    
    @staticmethod
    def create_fallback_file(file_path: str, error_msg: str) -> str:
        """Create a fallback error file when tool execution fails."""
        try:
            fallback_content = f"ERROR: Tool execution failed\nReason: {error_msg}\nTime: {datetime.now().isoformat()}"
            with open(file_path, 'w') as f:
                f.write(fallback_content)
            return f"Fallback file created: {file_path}"
        except Exception as e:
            return f"Failed to create fallback file: {str(e)}"
    
    @classmethod
    def execute_with_retry(cls, func, *args, **kwargs) -> Dict[str, Any]:
        """Execute a function with retry logic and error handling."""
        last_exception = None
        
        for attempt in range(MAX_RETRIES):
            try:
                result = func(*args, **kwargs)
                return {
                    'success': True,
                    'result': result,
                    'attempts': attempt + 1
                }
            except Exception as e:
                last_exception = e
                logger.warning(f"Attempt {attempt + 1} failed: {str(e)}")
                
                if not cls.should_retry(e, attempt + 1):
                    break
                
                delay = cls.get_delay(attempt)
                logger.info(f"Retrying in {delay:.1f} seconds...")
                time.sleep(delay)
        
        # All retries failed
        error_msg = f"{type(last_exception).__name__}: {str(last_exception)}"
        logger.error(f"All {MAX_RETRIES} attempts failed: {error_msg}")
        
        return {
            'success': False,
            'error': error_msg,
            'attempts': MAX_RETRIES,
            'fallback_suggestion': cls._get_recovery_suggestion(last_exception)
        }
    
    @staticmethod
    def _get_recovery_suggestion(exception: Exception) -> str:
        """Provide recovery suggestions based on error type."""
        if isinstance(exception, ImportError):
            return "Install missing dependency: pip install <package_name>"
        elif isinstance(exception, PermissionError):
            return "Check file permissions or close the file if it's open in another application"
        elif isinstance(exception, OSError):
            return "Check disk space and file system permissions"
        elif isinstance(exception, (ConnectionRefusedError, TimeoutError)):
            return "Ensure the external service is running and accessible"
        else:
            return "Review the error message and check tool configuration"

# =============================================================================
# TOOL EXECUTION FUNCTIONS
# =============================================================================

def execute_xls(payload: str) -> Dict[str, Any]:
    """Execute Excel tool to create a spreadsheet."""
    try:
        import pandas as pd
        
        parts = payload.split(':')
        filename = parts[0].strip()
        columns = [c.strip() for c in parts[1].split(',')]
        
        # Create sample data
        data = {col: [f"Sample {col} Data"] * 5 for col in columns}
        df = pd.DataFrame(data)
        
        # Get unique file path
        file_path = file_manager.get_unique_path(filename, '.xlsx')
        
        # Write to Excel
        df.to_excel(file_path, index=False)
        
        return {
            'tool': 'XLS',
            'file_path': file_path,
            'rows': len(df),
            'columns': len(columns)
        }
    except Exception as e:
        raise e

def execute_ppt(payload: str) -> Dict[str, Any]:
    """Execute PowerPoint tool to create a presentation."""
    try:
        from pptx import Presentation
        from pptx.util import Inches
        
        parts = payload.split(':')
        filename = parts[0].strip()
        slides_content = [s.strip() for s in parts[1].split('|')]
        
        # Create presentation
        prs = Presentation()
        
        # Add title slide
        if slides_content:
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = "Generated Presentation"
            if len(slides_content) > 0:
                slide.placeholders[1].text = slides_content[0]
        
        # Add content slides
        for content in slides_content[1:]:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = "Content Slide"
            slide.placeholders[1].text = content
        
        # Get unique file path
        file_path = file_manager.get_unique_path(filename, '.pptx')
        
        # Save presentation
        prs.save(file_path)
        
        return {
            'tool': 'PPT',
            'file_path': file_path,
            'slides': len(slides_content)
        }
    except Exception as e:
        raise e

def execute_pdf(payload: str) -> Dict[str, Any]:
    """Execute PDF tool to create a PDF document."""
    try:
        from fpdf import FPDF
        
        parts = payload.split(':', 1)
        filename = parts[0].strip()
        content = parts[1].strip()
        
        # Create PDF
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=12)
        
        # Add content (handle multi-line)
        for line in content.split('\n'):
            pdf.cell(200, 10, txt=line, ln=True)
        
        # Get unique file path
        file_path = file_manager.get_unique_path(filename, '.pdf')
        
        # Save PDF
        pdf.output(file_path)
        
        return {
            'tool': 'PDF',
            'file_path': file_path,
            'content_length': len(content)
        }
    except Exception as e:
        raise e

def execute_godot(payload: str) -> Dict[str, Any]:
    """Execute Godot tool to generate template files."""
    try:
        parts = payload.split(':')
        template_name = parts[0].strip()
        params = ':'.join(parts[1:]) if len(parts) > 1 else ""
        
        # Create a simple GDScript template
        script_content = f"""# Generated Godot Script
# Template: {template_name}
# Parameters: {params}
# Generated at: {datetime.now().isoformat()}

extends Node

func _ready():
    print("{template_name} initialized")
"""
        
        # Get unique file path
        file_path = file_manager.get_unique_path(template_name, '.gd')
        
        # Write script
        with open(file_path, 'w') as f:
            f.write(script_content)
        
        return {
            'tool': 'GODOT',
            'file_path': file_path,
            'template': template_name
        }
    except Exception as e:
        raise e

# =============================================================================
# SIGNAL DETECTION & ROUTING
# =============================================================================

# Signals where payload is a strict token (no spaces expected)
TOKEN_SIGNALS = {"XLS", "PPT", "GDT", "ERR"}

# Signals where payload is free text (spaces expected)
TEXT_SIGNALS = {"PDF", "WIKI", "WEB", "CHAT", "CODE", "PY", "CPP", "JS", "JSON"}

def extract_signals(text: str) -> List[Tuple[str, str]]:
    """Extract all signals from text in format @TOOL:payload"""
    results = []
    
    # Token-sensitive: stop at whitespace or @
    for m in re.finditer(r'@([A-Z]+):([^\s@]+)', text, re.IGNORECASE):
        if m.group(1).upper() in TOKEN_SIGNALS:
            results.append((m.group(1).upper(), m.group(2)))
    
    # Text-sensitive: capture to end of line or next signal
    for m in re.finditer(r'@([A-Z]+):(.+?)(?=\s*@[A-Z]+:|\s*$)', text, re.IGNORECASE):
        if m.group(1).upper() in TEXT_SIGNALS:
            results.append((m.group(1).upper(), m.group(2).strip()))
    
    return results

def has_signals(text: str) -> bool:
    """Check if text contains any tool signals."""
    return bool(extract_signals(text))

# =============================================================================
# MAIN EXECUTION FUNCTIONS
# =============================================================================

def execute_signal(signal_text: str) -> Dict[str, Any]:
    """
    Execute a single signal with full validation, caching, and error recovery.
    
    Args:
        signal_text: Signal in format "@TOOL:payload"
    
    Returns:
        Dictionary with execution result
    """
    # Extract tool and payload
    matches = re.match(r'@(\w+):(.+)', signal_text.strip())
    if not matches:
        return {
            'success': False,
            'error': 'Invalid signal format. Use @TOOL:payload'
        }
    
    tool_type = matches.group(1).upper()
    payload = matches.group(2)
    
    # 1. Check tool availability
    if not health_checker.is_tool_available(tool_type):
        return {
            'success': False,
            'error': f"Tool {tool_type} is not available. {health_checker.tool_status.get(tool_type, {}).get('reason', '')}"
        }
    
    # 2. Validate payload
    is_valid, error_msg = PayloadValidator.validate_payload(tool_type, payload)
    if not is_valid:
        return {
            'success': False,
            'error': f"Invalid payload: {error_msg}"
        }
    
    # 3. Check cache
    cached_result = tool_cache.get(tool_type, payload)
    if cached_result:
        logger.info(f"⚡ Cache hit for {tool_type}:{payload[:50]}...")
        cached_result['cached'] = True
        return cached_result
    
    # Map tool types to execution functions
    executors = {
        'XLS': execute_xls,
        'PPT': execute_ppt,
        'PDF': execute_pdf,
        'GODOT': execute_godot
    }
    
    executor = executors.get(tool_type)
    if not executor:
        return {
            'success': False,
            'error': f"Unknown tool type: {tool_type}"
        }
    
    # 5. Execute with error recovery
    logger.info(f"Executing {tool_type} tool with payload: {payload[:100]}...")
    result = ErrorRecoveryHandler.execute_with_retry(executor, payload)
    
    if result['success']:
        # Cache successful result
        tool_cache.set(tool_type, payload, result['result'])
        result['result']['cached'] = False
    
    return result

def execute_signals_parallel(text: str, max_workers: int = 5) -> List[Dict[str, Any]]:
    """
    Execute multiple signals in parallel.
    
    Args:
        text: Text containing multiple signals
        max_workers: Maximum number of concurrent workers
    
    Returns:
        List of execution results in original order
    """
    signals = extract_signals(text)
    
    if not signals:
        return [{'success': False, 'error': 'No signals found in text'}]
    
    results = [None] * len(signals)
    
    def execute_single(index: int, tool_type: str, payload: str):
        signal_str = f"@{tool_type}:{payload}"
        return index, execute_signal(signal_str)
    
    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        futures = [
            executor.submit(execute_single, idx, tool, payload)
            for idx, (tool, payload) in enumerate(signals)
        ]
        
        for future in as_completed(futures):
            idx, result = future.result()
            results[idx] = result
    
    return results

def get_available_tools() -> List[Dict[str, Any]]:
    """
    Return list of available tools with metadata.
    Only returns tools that pass health checks.
    """
    tools = []
    
    if health_checker.is_tool_available('XLS'):
        tools.append({
            'type': 'XLS',
            'signal': '@XLS:filename:col1,col2,col3',
            'description': 'Create Excel spreadsheet with specified columns',
            'example': '@XLS:Inventory:Item,Quantity,Price'
        })
    
    if health_checker.is_tool_available('PPT'):
        tools.append({
            'type': 'PPT',
            'signal': '@PPT:filename:slide1|slide2|slide3',
            'description': 'Create PowerPoint presentation with slides',
            'example': '@PPT:Report:Title|Overview|Details|Summary'
        })
    
    if health_checker.is_tool_available('PDF'):
        tools.append({
            'type': 'PDF',
            'signal': '@PDF:filename:content',
            'description': 'Create PDF document with text content',
            'example': '@PDF:Report:This is the report content...'
        })
    
    if health_checker.is_tool_available('GODOT'):
        tools.append({
            'type': 'GODOT',
            'signal': '@GODOT:template_name:params',
            'description': 'Generate Godot script template',
            'example': '@GODOT:PlayerController:move_speed=100,jump_force=10'
        })
    
    return tools

def get_tool_help() -> str:
    """Return formatted help text for available tools."""
    tools = get_available_tools()
    
    if not tools:
        return "❌ No tools are currently available. Check dependencies."
    
    lines = ["🛠️ Available Tools:", "=" * 40]
    for tool in tools:
        lines.append(f"\n{tool['type']}: {tool['description']}")
        lines.append(f"  Format: {tool['signal']}")
        lines.append(f"  Example: {tool['example']}")
    
    lines.append("\n" + "=" * 40)
    lines.append(health_checker.get_health_report())
    
    return "\n".join(lines)

# =============================================================================
# INITIALIZATION & TESTING
# =============================================================================

if __name__ == "__main__":
    print("🚀 COM Tool Harness Initialization")
    print("=" * 50)
    
    # Show health report
    print(health_checker.get_health_report())
    print()
    
    # Show available tools
    print(get_tool_help())
    print()
    
    # Test single signal execution
    print("\n🧪 Testing Single Signal Execution:")
    print("-" * 50)
    
    test_signal = "@XLS:TestReport:Item,Quantity,Price"
    print(f"Executing: {test_signal}")
    result = execute_signal(test_signal)
    print(f"Result: {result}")
    
    # Test cache
    print("\n🧪 Testing Cache (same signal):")
    print("-" * 50)
    result2 = execute_signal(test_signal)
    print(f"Result: {result2}")
    
    # Test parallel execution
    print("\n🧪 Testing Parallel Execution:")
    print("-" * 50)
    test_text = "@XLS:ParallelTest:A,B,C @GODOT:TestScript:param1=value1"
    print(f"Executing: {test_text}")
    parallel_results = execute_signals_parallel(test_text)
    for i, res in enumerate(parallel_results):
        print(f"Signal {i+1}: {res}")
    
    # Show final cache stats
    print("\n📊 Final Cache Statistics:")
    print(tool_cache.get_stats())
    
    print("\n✅ Tool harness ready for integration!")
