"""
COM SGMA LIGHT - PPT Tool
=========================
PowerPoint generation with comprehensive design system integration.
Integrated color systems: Tailwind CSS, Material Design, Ant Design, IBM Carbon, 
GitHub Primer, Radix UI, Chakra UI, Atlassian, and Open Color.
"""
# Requires: pip install python-pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ============================================================================
# COMPREHENSIVE DESIGN TOKENS - Integrated from Multiple Design Systems
# Matches PDF tool for cross-format consistency
# ============================================================================

# Brand Core - COM SGMA LIGHT
BRAND_COLORS = {
    "primary": "#1a5f3c",      # Emerald green - primary actions, headers
    "secondary": "#0d3d26",    # Dark emerald - accents, borders
    "accent": "#ffb1ee",       # Soft pink - highlights, CTAs
    "background": "#ffffff",   # White - page background
    "surface": "#f8f9fa",      # Light gray - content blocks
    "text_primary": "#1a1a1a", # Near black - main text
    "text_secondary": "#666666", # Gray - secondary text
    "border": "#e0e0e0",       # Light border
}

# Tailwind CSS Colors - Industry Standard (120K+ stars)
TAILWIND_COLORS = {
    "red": {"50": "#fef2f2", "100": "#fee2e2", "200": "#fecaca", "300": "#fca5a5", "400": "#f87171", "500": "#ef4444", "600": "#dc2626", "700": "#b91c1c", "800": "#991b1b", "900": "#7f1d1d"},
    "orange": {"50": "#fff7ed", "100": "#ffedd5", "200": "#fed7aa", "300": "#fdba74", "400": "#fb923c", "500": "#f97316", "600": "#ea580c", "700": "#c2410c", "800": "#9a3412", "900": "#7c2d12"},
    "amber": {"50": "#fffbeb", "100": "#fef3c7", "200": "#fde68a", "300": "#fcd34d", "400": "#fbbf24", "500": "#f59e0b", "600": "#d97706", "700": "#b45309", "800": "#92400e", "900": "#78350f"},
    "yellow": {"50": "#fefce8", "100": "#fef9c3", "200": "#fef08a", "300": "#fde047", "400": "#facc15", "500": "#eab308", "600": "#ca8a04", "700": "#a16207", "800": "#854d0e", "900": "#713f12"},
    "green": {"50": "#f0fdf4", "100": "#dcfce7", "200": "#bbf7d0", "300": "#86efac", "400": "#4ade80", "500": "#22c55e", "600": "#16a34a", "700": "#15803d", "800": "#166534", "900": "#14532d"},
    "emerald": {"50": "#ecfdf5", "100": "#d1fae5", "200": "#a7f3d0", "300": "#6ee7b7", "400": "#34d399", "500": "#10b981", "600": "#059669", "700": "#047857", "800": "#065f46", "900": "#064e3b"},
    "teal": {"50": "#f0fdfa", "100": "#ccfbf1", "200": "#99f6e4", "300": "#5eead4", "400": "#2dd4bf", "500": "#14b8a6", "600": "#0d9488", "700": "#0f766e", "800": "#115e59", "900": "#134e4a"},
    "cyan": {"50": "#ecfeff", "100": "#cffafe", "200": "#a5f3fc", "300": "#67e8f9", "400": "#22d3ee", "500": "#06b6d4", "600": "#0891b2", "700": "#0e7490", "800": "#155e75", "900": "#164e63"},
    "sky": {"50": "#f0f9ff", "100": "#e0f2fe", "200": "#bae6fd", "300": "#7dd3fc", "400": "#38bdf8", "500": "#0ea5e9", "600": "#0284c7", "700": "#0369a1", "800": "#075985", "900": "#0c4a6e"},
    "blue": {"50": "#eff6ff", "100": "#dbeafe", "200": "#bfdbfe", "300": "#93c5fd", "400": "#60a5fa", "500": "#3b82f6", "600": "#2563eb", "700": "#1d4ed8", "800": "#1e40af", "900": "#1e3a8a"},
    "indigo": {"50": "#eef2ff", "100": "#e0e7ff", "200": "#c7d2fe", "300": "#a5b4fc", "400": "#818cf8", "500": "#6366f1", "600": "#4f46e5", "700": "#4338ca", "800": "#3730a3", "900": "#312e81"},
    "violet": {"50": "#f5f3ff", "100": "#ede9fe", "200": "#ddd6fe", "300": "#c4b5fd", "400": "#a78bfa", "500": "#8b5cf6", "600": "#7c3aed", "700": "#6d28d9", "800": "#5b21b6", "900": "#4c1d95"},
    "purple": {"50": "#faf5ff", "100": "#f3e8ff", "200": "#e9d5ff", "300": "#d8b4fe", "400": "#c084fc", "500": "#a855f7", "600": "#9333ea", "700": "#7e22ce", "800": "#6b21a8", "900": "#581c87"},
    "fuchsia": {"50": "#fdf4ff", "100": "#fae8ff", "200": "#f5d0fe", "300": "#f0abfc", "400": "#e879f9", "500": "#d946ef", "600": "#c026d3", "700": "#a21caf", "800": "#86198f", "900": "#701a75"},
    "pink": {"50": "#fdf2f8", "100": "#fce7f3", "200": "#fbcfe8", "300": "#f9a8d4", "400": "#f472b6", "500": "#ec4899", "600": "#db2777", "700": "#be185d", "800": "#9d174d", "900": "#831843"},
    "rose": {"50": "#fff1f2", "100": "#ffe4e6", "200": "#fecdd3", "300": "#fda4af", "400": "#fb7185", "500": "#f43f5e", "600": "#e11d48", "700": "#be123c", "800": "#9f1239", "900": "#881337"},
    "slate": {"50": "#f8fafc", "100": "#f1f5f9", "200": "#e2e8f0", "300": "#cbd5e1", "400": "#94a3b8", "500": "#64748b", "600": "#475569", "700": "#334155", "800": "#1e293b", "900": "#0f172a"},
    "gray": {"50": "#f9fafb", "100": "#f3f4f6", "200": "#e5e7eb", "300": "#d1d5db", "400": "#9ca3af", "500": "#6b7280", "600": "#4b5563", "700": "#374151", "800": "#1f2937", "900": "#111827"},
    "zinc": {"50": "#fafafa", "100": "#f4f4f5", "200": "#e4e4e7", "300": "#d4d4d8", "400": "#a1a1aa", "500": "#71717a", "600": "#52525b", "700": "#3f3f46", "800": "#27272a", "900": "#18181b"},
    "neutral": {"50": "#fafafa", "100": "#f5f5f5", "200": "#e5e5e5", "300": "#d4d4d4", "400": "#a3a3a3", "500": "#737373", "600": "#525252", "700": "#404040", "800": "#262626", "900": "#171717"},
    "stone": {"50": "#fafaf9", "100": "#f5f5f4", "200": "#e7e5e4", "300": "#d6d3d1", "400": "#a8a29e", "500": "#78716c", "600": "#57534e", "700": "#44403c", "800": "#292524", "900": "#1c1917"},
}

# Material Design Colors - Google's Accessibility-Tested System (40K+ stars)
MATERIAL_COLORS = {
    "red": {"50": "#ffebee", "100": "#ffcdd2", "200": "#ef9a9a", "300": "#e57373", "400": "#ef5350", "500": "#f44336", "600": "#e53935", "700": "#d32f2f", "800": "#c62828", "900": "#b71c1c", "A100": "#ff8a80", "A200": "#ff5252", "A400": "#ff1744", "A700": "#d50000"},
    "pink": {"50": "#fce4ec", "100": "#f8bbd0", "200": "#f48fb1", "300": "#f06292", "400": "#ec407a", "500": "#e91e63", "600": "#d81b60", "700": "#c2185b", "800": "#ad1457", "900": "#880e4f", "A100": "#ff80ab", "A200": "#ff4081", "A400": "#f50057", "A700": "#c51162"},
    "purple": {"50": "#f3e5f5", "100": "#e1bee7", "200": "#ce93d8", "300": "#ba68c8", "400": "#ab47bc", "500": "#9c27b0", "600": "#8e24aa", "700": "#7b1fa2", "800": "#6a1b9a", "900": "#4a148c", "A100": "#ea80fc", "A200": "#e040fb", "A400": "#d500f9", "A700": "#aa00ff"},
    "blue": {"50": "#e3f2fd", "100": "#bbdefb", "200": "#90caf9", "300": "#64b5f6", "400": "#42a5f5", "500": "#2196f3", "600": "#1e88e5", "700": "#1976d2", "800": "#1565c0", "900": "#0d47a1", "A100": "#82b1ff", "A200": "#448aff", "A400": "#2979ff", "A700": "#2962ff"},
    "cyan": {"50": "#e0f7fa", "100": "#b2ebf2", "200": "#80deea", "300": "#4dd0e1", "400": "#26c6da", "500": "#00bcd4", "600": "#00acc1", "700": "#0097a7", "800": "#00838f", "900": "#006064", "A100": "#84ffff", "A200": "#18ffff", "A400": "#00e5ff", "A700": "#00b8d4"},
    "teal": {"50": "#e0f2f1", "100": "#b2dfdb", "200": "#80cbc4", "300": "#4db6ac", "400": "#26a69a", "500": "#009688", "600": "#00897b", "700": "#00796b", "800": "#00695c", "900": "#004d40", "A100": "#a7ffeb", "A200": "#64ffda", "A400": "#1de9b6", "A700": "#00bfa5"},
    "green": {"50": "#e8f5e9", "100": "#c8e6c9", "200": "#a5d6a7", "300": "#81c784", "400": "#66bb6a", "500": "#4caf50", "600": "#43a047", "700": "#388e3c", "800": "#2e7d32", "900": "#1b5e20", "A100": "#b9f6ca", "A200": "#69f0ae", "A400": "#00e676", "A700": "#00c853"},
    "yellow": {"50": "#fffde7", "100": "#fff9c4", "200": "#fff59d", "300": "#fff176", "400": "#ffee58", "500": "#ffeb3b", "600": "#fdd835", "700": "#fbc02d", "800": "#f9a825", "900": "#f57f17", "A100": "#ffff8d", "A200": "#ffff00", "A400": "#ffea00", "A700": "#ffd600"},
    "orange": {"50": "#fff3e0", "100": "#ffe0b2", "200": "#ffcc80", "300": "#ffb74d", "400": "#ffa726", "500": "#ff9800", "600": "#fb8c00", "700": "#f57c00", "800": "#ef6c00", "900": "#e65100", "A100": "#ffd180", "A200": "#ffab40", "A400": "#ff9100", "A700": "#ff6d00"},
    "grey": {"50": "#fafafa", "100": "#f5f5f5", "200": "#eeeeee", "300": "#e0e0e0", "400": "#bdbdbd", "500": "#9e9e9e", "600": "#757575", "700": "#616161", "800": "#424242", "900": "#212121"},
}

# Ant Design Colors - Asian Market Optimized (90K+ stars)
ANT_DESIGN_COLORS = {
    "blue": {"1": "#e6f4ff", "2": "#bae0ff", "3": "#91caff", "4": "#69b1ff", "5": "#4096ff", "6": "#1677ff", "7": "#0958d9", "8": "#003eb3", "9": "#002c8c", "10": "#001d66"},
    "purple": {"1": "#f9f0ff", "2": "#efdbff", "3": "#d3adf7", "4": "#b37feb", "5": "#9254de", "6": "#722ed1", "7": "#531dab", "8": "#391085", "9": "#22075e", "10": "#120338"},
    "cyan": {"1": "#e6fffb", "2": "#b5f5ec", "3": "#87e8de", "4": "#5cdbd3", "5": "#22c1c3", "6": "#00a4a4", "7": "#00898b", "8": "#006b75", "9": "#00505b", "10": "#003442"},
    "green": {"1": "#f6ffed", "2": "#d9f7be", "3": "#b7eb8f", "4": "#95de64", "5": "#73d13d", "6": "#52c41a", "7": "#389e0d", "8": "#237804", "9": "#135200", "10": "#092b00"},
    "red": {"1": "#fff1f0", "2": "#ffccc7", "3": "#ffa39e", "4": "#ff7875", "5": "#ff4d4f", "6": "#f5222d", "7": "#cf1322", "8": "#a8071a", "9": "#7a0413", "10": "#5c0011"},
    "orange": {"1": "#fff7e6", "2": "#ffe7ba", "3": "#ffd591", "4": "#ffc069", "5": "#ffa940", "6": "#fa8c16", "7": "#d46b08", "8": "#ad4e00", "9": "#873800", "10": "#612500"},
    "yellow": {"1": "#feffe6", "2": "#ffffb8", "3": "#fffb8f", "4": "#fff566", "5": "#ffec3d", "6": "#fadb14", "7": "#d4b106", "8": "#ad8b00", "9": "#876800", "10": "#614700"},
    "gold": {"1": "#fffbe6", "2": "#fff1b8", "3": "#ffe58f", "4": "#ffd666", "5": "#ffc53d", "6": "#faad14", "7": "#d48806", "8": "#ad6800", "9": "#874d00", "10": "#613400"},
}

# IBM Carbon Colors - Enterprise Professional (8K+ stars)
IBM_CARBON_COLORS = {
    "blue": {"10": "#edf5ff", "20": "#d0e2ff", "30": "#a6c8ff", "40": "#78a9ff", "50": "#4589ff", "60": "#0f62fe", "70": "#0043ce", "80": "#002d9c", "90": "#001d6c", "100": "#001141"},
    "cyan": {"10": "#e5f9ff", "20": "#baefff", "30": "#82ddff", "40": "#33cbff", "50": "#00b3ff", "60": "#0092ff", "70": "#0072c3", "80": "#005387", "90": "#003552", "100": "#001a29"},
    "gray": {"10": "#f6f6f6", "20": "#e5e5e5", "30": "#c6c6c6", "40": "#a8a8a8", "50": "#8d8d8d", "60": "#6f6f6f", "70": "#525252", "80": "#393939", "90": "#262626", "100": "#161616"},
    "green": {"10": "#defbe6", "20": "#a7f4c4", "30": "#6ee3a0", "40": "#37d07b", "50": "#19c060", "60": "#0ba947", "70": "#078a38", "80": "#056b2b", "90": "#034d1e", "100": "#013013"},
    "red": {"10": "#fff1f1", "20": "#ffd7d9", "30": "#ffb3b8", "40": "#ff7d7f", "50": "#ff4d4f", "60": "#da1e28", "70": "#a2191f", "80": "#750e13", "90": "#520408", "100": "#2d0304"},
}

# GitHub Primer Colors - Data Visualization Optimized (3K+ stars)
GITHUB_PRIMER_COLORS = {
    "coral": {"1": "#fff5f1", "2": "#ffe6d5", "3": "#ffd0b5", "4": "#ffb48c", "5": "#ff925f", "6": "#ff6f2e", "7": "#e85a1a", "8": "#bc4413", "9": "#94320e", "10": "#75250a"},
    "gray": {"1": "#f6f8fa", "2": "#eaeef2", "3": "#d0d7de", "4": "#afb8c1", "5": "#8c959f", "6": "#6e7781", "7": "#57606a", "8": "#424a53", "9": "#32383f", "10": "#25292e"},
    "blue": {"1": "#ddf4ff", "2": "#b6e3ff", "3": "#80ccff", "4": "#54aeff", "5": "#218bff", "6": "#0969da", "7": "#0550ae", "8": "#033d8b", "9": "#0a3069", "10": "#001a47"},
    "green": {"1": "#dafbe1", "2": "#aceebb", "3": "#6fdd8b", "4": "#4ac26b", "5": "#2da44e", "6": "#1a7f37", "7": "#116329", "8": "#044f1e", "9": "#003d16", "10": "#002d11"},
    "orange": {"1": "#fff1e5", "2": "#ffd8b5", "3": "#ffb775", "4": "#ff8d32", "5": "#fd7504", "6": "#df5c02", "7": "#b54702", "8": "#8d3504", "9": "#692107", "10": "#4f170b"},
    "pink": {"1": "#ffeef8", "2": "#fad1df", "3": "#f7a8c7", "4": "#f575ad", "5": "#ef4e8a", "6": "#d0356e", "7": "#a82858", "8": "#821c45", "9": "#5e1433", "10": "#420e25"},
    "purple": {"1": "#f7f1ff", "2": "#e9d7fe", "3": "#d8b9fe", "4": "#c297ff", "5": "#a475f9", "6": "#8957e5", "7": "#6e40c9", "8": "#55309e", "9": "#3f2374", "10": "#2d1a52"},
    "red": {"1": "#ffebe9", "2": "#ffcecb", "3": "#ffaba8", "4": "#ff8182", "5": "#ff4f4a", "6": "#da3633", "7": "#b62324", "8": "#8e1519", "9": "#670a0e", "10": "#4b090b"},
}

# Radix UI Colors - 30+ Families with Dark Mode (5K+ stars)
RADIX_COLORS = {
    "tomato": {"1": "#fff5f3", "2": "#ffe2db", "3": "#ffc6b8", "4": "#ffa58f", "5": "#ff8066", "6": "#ff5940", "7": "#fa3d26", "8": "#d92e1a", "9": "#c62818", "10": "#ab2215"},
    "ruby": {"1": "#fff5f6", "2": "#ffe2e6", "3": "#ffc6cf", "4": "#ffa5b4", "5": "#ff8096", "6": "#ff5976", "7": "#fa3d56", "8": "#d92e3f", "9": "#c62838", "10": "#ab222f"},
    "crimson": {"1": "#fff5f7", "2": "#ffe2e7", "3": "#ffc6d0", "4": "#ffa5b6", "5": "#ff8098", "6": "#ff5976", "7": "#fa3d56", "8": "#d92e41", "9": "#c6283a", "10": "#ab2231"},
    "plum": {"1": "#faf5fa", "2": "#f5e2f5", "3": "#edc6ed", "4": "#e2a5e2", "5": "#d680d6", "6": "#c659c6", "7": "#b53db5", "8": "#9e2e9e", "9": "#902890", "10": "#7c227c"},
    "mint": {"1": "#f5fbf7", "2": "#e2f5e7", "3": "#c6ebcf", "4": "#a5deb3", "5": "#80d093", "6": "#59be71", "7": "#3da854", "8": "#2e8f42", "9": "#28823c", "10": "#226f33"},
}

# Chakra UI Colors - Simple Naming Convention (35K+ stars)
CHAKRA_COLORS = {
    "brand": {"50": "#e3f2fd", "100": "#bbdefb", "200": "#90caf9", "300": "#64b5f6", "400": "#42a5f5", "500": "#3182ce", "600": "#2b6cb0", "700": "#2c5282", "800": "#2a4365", "900": "#1a365d"},
    "linkedin": {"50": "#e8f4fd", "100": "#c8e6f9", "200": "#a4d4f4", "300": "#7ec1ef", "400": "#5bafea", "500": "#3b99e3", "600": "#2f7dbd", "700": "#256296", "800": "#1c4a70", "900": "#123149"},
    "facebook": {"50": "#e8f4fd", "100": "#c8e6f9", "200": "#a4d4f4", "300": "#7ec1ef", "400": "#5bafea", "500": "#3b5998", "600": "#2f477a", "700": "#253861", "800": "#1c2b49", "900": "#121c30"},
    "whatsapp": {"50": "#e8f5e9", "100": "#c8e6c9", "200": "#a5d6a7", "300": "#81c784", "400": "#66bb6a", "500": "#25D366", "600": "#1fa851", "700": "#188540", "800": "#126430", "900": "#0b411f"},
}

# Atlassian Design Colors - Business Presentation Optimized (5K+ stars)
ATLASSIAN_COLORS = {
    "blue": {"N10": "#deebff", "N20": "#b3d4ff", "N30": "#80b8ff", "N40": "#4d9eff", "N50": "#2684ff", "N60": "#0065ff", "N70": "#0052cc", "N80": "#003d99", "N90": "#002a66", "N100": "#001a40"},
    "teal": {"N10": "#e3fcf7", "N20": "#b3f5ea", "N30": "#80ebe0", "N40": "#4dded3", "N50": "#26cfbf", "N60": "#00bfa5", "N70": "#00998a", "N80": "#007366", "N90": "#004d44", "N100": "#002622"},
    "green": {"N10": "#e3fcef", "N20": "#b3f5d4", "N30": "#80ebbb", "N40": "#4dde9f", "N50": "#26cf82", "N60": "#00bfa5", "N70": "#00998a", "N80": "#007366", "N90": "#004d44", "N100": "#002622"},
    "red": {"N10": "#ffebe6", "N20": "#ffcdc7", "N30": "#ffab9e", "N40": "#ff806e", "N50": "#ff563d", "N60": "#ff3621", "N70": "#de350b", "N80": "#bf2600", "N90": "#8f1a00", "N100": "#5c1100"},
    "purple": {"N10": "#eae6ff", "N20": "#cdc7ff", "N30": "#ab9eff", "N40": "#806eff", "N50": "#563dff", "N60": "#403294", "N70": "#362b7a", "N80": "#2a215c", "N90": "#1d1740", "N100": "#110d26"},
}

# Open Color - Scandinavian Design (Original Request)
OPEN_COLOR = {
    "gray": ["#f8f9fa", "#f1f3f5", "#e9ecef", "#dee2e6", "#ced4da", "#adb5bd", "#868e96", "#495057", "#343a40", "#212529"],
    "red": ["#fff5f5", "#ffe3e3", "#ffc9c9", "#ffa8a8", "#ff8787", "#ff6b6b", "#fa5252", "#f03e3e", "#e03131", "#c92a2a"],
    "pink": ["#fff0f6", "#ffdeeb", "#fcc2d7", "#faa2c1", "#f783ac", "#f06595", "#e64980", "#d6336c", "#c2255c", "#a61e4d"],
    "grape": ["#f8f0fc", "#f3d9fa", "#eebefa", "#e599f7", "#da77f2", "#cc5de8", "#be4bdb", "#ae3ec9", "#9c36b5", "#862e9c"],
    "violet": ["#f3f0ff", "#e5dbff", "#d0bfff", "#b197fc", "#9775fa", "#845ef7", "#7950f2", "#7048e8", "#6741d9", "#5f3dc4"],
    "indigo": ["#edf2ff", "#dbe4ff", "#bac8ff", "#91a7ff", "#748ffc", "#5c7cfa", "#4c6ef5", "#4263eb", "#3b5bdb", "#364fc7"],
    "blue": ["#e7f5ff", "#d0ebff", "#a5d8ff", "#74c0fc", "#4dabf7", "#339af0", "#228be6", "#1c7ed6", "#1971c2", "#1864ab"],
    "cyan": ["#e3fafc", "#c5f6fa", "#99e9f2", "#66d9e8", "#3bc9db", "#22b8cf", "#15aabf", "#1098ad", "#0c8599", "#0b7285"],
    "teal": ["#e6fcf5", "#c3fae8", "#96f2d7", "#63e6be", "#38d9a9", "#20c997", "#12b886", "#0ca678", "#099268", "#087f5b"],
    "green": ["#ebfbee", "#d3f9d8", "#b2f2bb", "#8ce99a", "#69db7c", "#51cf66", "#40c057", "#37b24d", "#2f9e44", "#2b8a3e"],
    "lime": ["#f4fce3", "#e9fac8", "#d8f5a2", "#c0eb75", "#a9e34b", "#94d82d", "#82c91e", "#74b816", "#66a80f", "#5c940d"],
    "yellow": ["#fff9db", "#fff3bf", "#ffec99", "#ffe066", "#ffd43b", "#fcc419", "#fab005", "#f59f00", "#f08c00", "#e67700"],
    "orange": ["#fff4e6", "#ffe8cc", "#ffd8a8", "#ffc078", "#ffa94d", "#ff922b", "#fd7e14", "#f76707", "#e8590c", "#d9480f"],
}

# Consolidated Design Tokens - Primary selections from each system
DESIGN_TOKENS = {
    "colors": {
        # Brand Core
        **BRAND_COLORS,
        # Extended palette references (use via get_color function)
        "_systems": {
            "tailwind": TAILWIND_COLORS,
            "material": MATERIAL_COLORS,
            "ant_design": ANT_DESIGN_COLORS,
            "ibm_carbon": IBM_CARBON_COLORS,
            "github_primer": GITHUB_PRIMER_COLORS,
            "radix": RADIX_COLORS,
            "chakra": CHAKRA_COLORS,
            "atlassian": ATLASSIAN_COLORS,
            "open_color": OPEN_COLOR,
        }
    },
    "typography": {
        "heading_font": "Arial",    # Clean sans-serif for headings
        "body_font": "Arial",       # Consistent body font
    }
}


def get_color(system: str, color_name: str, shade: str = "500") -> str:
    """
    Get color from any integrated design system.
    
    Args:
        system: Design system name (tailwind, material, ant_design, ibm_carbon, 
                github_primer, radix, chakra, atlassian, open_color)
        color_name: Color family name (blue, red, green, etc.)
        shade: Shade level (varies by system: 50-900, 1-10, N10-N100, etc.)
    
    Returns:
        Hex color string or None if not found
    """
    systems_map = DESIGN_TOKENS["colors"]["_systems"]
    if system not in systems_map:
        return BRAND_COLORS["primary"]
    
    palette = systems_map[system]
    if color_name not in palette:
        return BRAND_COLORS["primary"]
    
    shades = palette[color_name]
    if isinstance(shades, dict):
        return shades.get(str(shade), BRAND_COLORS["primary"])
    elif isinstance(shades, list):
        try:
            idx = int(shade) if shade.isdigit() else 5
            return shades[min(idx, len(shades)-1)]
        except:
            return shades[5] if len(shades) > 5 else shades[0]
    
    return BRAND_COLORS["primary"]


def run(payload: str, template: str = "default") -> str:
    """
    Payload format: filename:slide1|slide2|slide3
    Example: Deck:Introduction|Data Results|Conclusion
    
    Templates available:
    - default: Clean professional slides
    - corporate: Business presentation with header bar
    - modern: Bold accent colors, minimal layout
    - dark: Dark theme for presentations
    
    Args:
        payload: String containing filename and slide titles (pipe-separated)
        template: Template style to apply
    
    Returns:
        Status string
    """
    try:
        parts = payload.split(":", 1)
        filename = parts[0].strip() if parts[0] else "presentation"
        slides_raw = parts[1].strip() if len(parts) > 1 else "Slide 1"
        slide_titles = [s.strip() for s in slides_raw.split("|")]

        prs = Presentation()
        
        # Apply template styling
        if template == "corporate":
            _apply_corporate_template(prs, slide_titles, DESIGN_TOKENS)
        elif template == "modern":
            _apply_modern_template(prs, slide_titles, DESIGN_TOKENS)
        elif template == "dark":
            _apply_dark_template(prs, slide_titles, DESIGN_TOKENS)
        else:  # default
            _apply_default_template(prs, slide_titles, DESIGN_TOKENS)

        path = f"{filename}.pptx"
        prs.save(path)
        return f"✅ PPT created: {path} ({template} template) | {len(slide_titles)} slides"

    except Exception as e:
        return f"❌ PPT failed: {str(e)}"


def _apply_default_template(prs: Presentation, slide_titles: list, tokens: dict) -> None:
    """Clean professional slides with subtle branding."""
    colors = tokens["colors"]
    layout = prs.slide_layouts[1]  # title + content
    
    for title_text in slide_titles:
        slide = prs.slides.add_slide(layout)
        
        # Title styling
        title_shape = slide.shapes.title
        title_shape.text = title_text
        title_tf = title_shape.text_frame.paragraphs[0]
        title_tf.font.size = Pt(32)
        title_tf.font.color.rgb = RGBColor(*_hex_to_rgb(colors["primary"]))
        title_tf.font.bold = True
        
        # Subtle accent line below title
        left = Inches(0.5)
        top = Inches(1.8)
        width = Inches(2)
        height = Pt(3)
        accent_line = slide.shapes.add_shape(
            1, left, top, width, height  # msoShapeRectangle
        )
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = RGBColor(*_hex_to_rgb(colors["accent"]))
        accent_line.line.fill.background()


def _apply_corporate_template(prs: Presentation, slide_titles: list, tokens: dict) -> None:
    """Business presentation with header bar and formal layout."""
    colors = tokens["colors"]
    layout = prs.slide_layouts[1]
    
    for title_text in slide_titles:
        slide = prs.slides.add_slide(layout)
        
        # Header bar at top
        left = Inches(0)
        top = Inches(0)
        width = Inches(10)
        height = Inches(0.6)
        header = slide.shapes.add_shape(1, left, top, width, height)
        header.fill.solid()
        header.fill.fore_color.rgb = RGBColor(*_hex_to_rgb(colors["primary"]))
        header.line.fill.background()
        
        # Title (positioned below header)
        title_shape = slide.shapes.title
        title_shape.top = Inches(0.8)
        title_shape.text = title_text
        title_tf = title_shape.text_frame.paragraphs[0]
        title_tf.font.size = Pt(28)
        title_tf.font.color.rgb = RGBColor(*_hex_to_rgb(colors["text_primary"]))
        title_tf.font.bold = True
        
        # Footer with branding
        left = Inches(0.5)
        top = Inches(7)
        footer = slide.shapes.add_textbox(left, top, Inches(9), Inches(0.3))
        footer_tf = footer.text_frame
        footer_p = footer_tf.paragraphs[0]
        footer_p.text = "COM SGMA LIGHT | CONFIDENTIAL"
        footer_p.font.size = Pt(9)
        footer_p.font.color.rgb = RGBColor(*_hex_to_rgb(colors["text_secondary"]))
        footer_p.alignment = PP_ALIGN.CENTER


def _apply_modern_template(prs: Presentation, slide_titles: list, tokens: dict) -> None:
    """Bold accent colors, minimal layout with large typography."""
    colors = tokens["colors"]
    layout = prs.slide_layouts[6]  # blank layout for custom design
    
    for title_text in slide_titles:
        slide = prs.slides.add_slide(layout)
        
        # Accent block on left side
        left = Inches(0)
        top = Inches(0)
        width = Inches(0.4)
        height = Inches(7.5)
        accent_block = slide.shapes.add_shape(1, left, top, width, height)
        accent_block.fill.solid()
        accent_block.fill.fore_color.rgb = RGBColor(*_hex_to_rgb(colors["accent"]))
        accent_block.line.fill.background()
        
        # Large title
        left = Inches(1)
        top = Inches(2.5)
        width = Inches(8)
        height = Inches(2)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_tf = title_box.text_frame
        title_p = title_tf.paragraphs[0]
        title_p.text = title_text
        title_p.font.size = Pt(44)
        title_p.font.color.rgb = RGBColor(*_hex_to_rgb(colors["text_primary"]))
        title_p.font.bold = True


def _apply_dark_template(prs: Presentation, slide_titles: list, tokens: dict) -> None:
    """Dark theme for presentations - high contrast, cinematic."""
    # Dark theme specific colors
    dark_bg = "#1a1a1a"
    dark_surface = "#2d2d2d"
    light_text = "#ffffff"
    accent = "#ffb1ee"
    
    layout = prs.slide_layouts[6]  # blank layout
    
    for title_text in slide_titles:
        slide = prs.slides.add_slide(layout)
        
        # Dark background
        bg = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(10), Inches(7.5)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(*_hex_to_rgb(dark_bg))
        bg.line.fill.background()
        
        # Accent gradient effect (multiple rectangles)
        for i in range(3):
            left = Inches(7 + i * 0.3)
            top = Inches(0.5 + i * 0.2)
            width = Inches(2.5)
            height = Inches(6)
            accent_shape = slide.shapes.add_shape(1, left, top, width, height)
            accent_shape.fill.solid()
            accent_shape.fill.fore_color.rgb = RGBColor(*_hex_to_rgb(accent))
            accent_shape.fill.transparency = 0.7 - (i * 0.2)
            accent_shape.line.fill.background()
        
        # White title
        left = Inches(0.8)
        top = Inches(2.5)
        width = Inches(6)
        height = Inches(2)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_tf = title_box.text_frame
        title_p = title_tf.paragraphs[0]
        title_p.text = title_text
        title_p.font.size = Pt(40)
        title_p.font.color.rgb = RGBColor(*_hex_to_rgb(light_text))
        title_p.font.bold = True


def _hex_to_rgb(hex_color: str) -> tuple:
    """Convert hex color to RGB tuple."""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
